import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def apply_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_title(slide, step_num, title_text):
    # 用一個小巧的藍色矩形做裝飾
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.6), Inches(0.15), Inches(0.6))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(37, 99, 235) # 商務藍 #2563EB
    rect.line.fill.background() # 無邊框
    
    # 標題文字框
    txBox = slide.shapes.add_textbox(Inches(1.1), Inches(0.5), Inches(11), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = f"{step_num}  {title_text}"
    p.font.name = 'Microsoft JhengHei'
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42) # Slate 900

def create_card(slide, left, top, width, height, bg_color, border_color=None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
    else:
        card.line.fill.background()
    return card

def add_bullet_points(slide, left, top, width, height, points, title=None, title_size=18, text_size=14, is_dark_bg=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.1)
    
    text_color = RGBColor(255, 255, 255) if is_dark_bg else RGBColor(15, 23, 42)
    sub_color = RGBColor(226, 232, 240) if is_dark_bg else RGBColor(71, 85, 105)
    
    first = True
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(title_size)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.space_after = Pt(10)
        first = False
        
    for pt in points:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
            
        p.text = pt
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(text_size)
        p.font.color.rgb = text_color
        p.space_after = Pt(6)
        
        # 如果是項目符號
        if pt.startswith('• '):
            p.text = pt # 保留項目符號作為文字
            # 調整行距以防重疊
            p.line_spacing = 1.15

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # 空白投影片
    
    # ----------------------------------------------------
    # Slide 1: 封面
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, RGBColor(248, 250, 252))
    
    # 左側商務藍色條
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(37, 99, 235)
    rect.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.5), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = 0
    
    p1 = tf.paragraphs[0]
    p1.text = "AIoT 智慧安全監控與主動防禦平台"
    p1.font.name = 'Microsoft JhengHei'
    p1.font.size = Pt(42)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(15, 23, 42)
    p1.space_after = Pt(14)
    
    p2 = tf.add_paragraph()
    p2.text = "整合 YOLO26n 影像辨識與 Gemini 3.5 AI 智慧稽核的工安主動防禦系統"
    p2.font.name = 'Microsoft JhengHei'
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(71, 85, 105)
    p2.space_after = Pt(40)
    
    p3 = tf.add_paragraph()
    p3.text = "專案成果簡報  |  2026年6月"
    p3.font.name = 'Microsoft JhengHei'
    p3.font.size = Pt(13)
    p3.font.color.rgb = RGBColor(148, 163, 184)
    
    # ----------------------------------------------------
    # Slide 2: 專案背景與工安挑戰
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, RGBColor(248, 250, 252))
    create_title(slide, "01", "專案背景與工安挑戰")
    
    # 左欄
    create_card(slide, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0), RGBColor(255, 255, 255), RGBColor(239, 68, 68))
    add_bullet_points(slide, Inches(1.1), Inches(1.9), Inches(5.0), Inches(4.4), [
        "• 人工巡檢存在盲點：廠區範圍廣大，安全管理人員無法做到 24 小時無死角監視，易生疏漏。",
        "• 意外發生反應遲緩：傳統監控僅有錄影功能，發現違規到實際通報存在時間差，無法即時阻止意外。",
        "• 歷史數據難以追蹤：工安違規紀錄零散，難以進行數位化統計、趨勢分析與系統性預防改善。"
    ], title="傳統工安痛點", title_size=20, text_size=15)
    
    # 右欄
    create_card(slide, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.0), RGBColor(255, 255, 255), RGBColor(37, 99, 235))
    add_bullet_points(slide, Inches(7.2), Inches(1.9), Inches(5.0), Inches(4.4), [
        "• 邊緣端 AI 即時影像辨識：透過邊緣運算執行 YOLO 模型，可實現毫秒級的高精度穿戴辨識。",
        "• 大語言模型智慧稽核：藉由 LLM (Gemini 3.5) 的理解力，快速分析歷史違規紀錄，找出管理盲點。",
        "• 即時聯網警報通報：結合 Discord Webhook，將違規畫面與通知即時推播給管理人員，落實主動防禦。"
    ], title="科技化轉型契機", title_size=20, text_size=15)
    
    # ----------------------------------------------------
    # Slide 3: 系統核心功能特點
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, RGBColor(248, 250, 252))
    create_title(slide, "02", "系統核心功能特點")
    
    col_width = Inches(3.6)
    gap = Inches(0.4)
    start_left = Inches(0.8)
    top = Inches(1.8)
    height = Inches(4.6)
    
    # 卡片 1 (YOLO)
    create_card(slide, start_left, top, col_width, height, RGBColor(255, 255, 255), RGBColor(16, 185, 129))
    add_bullet_points(slide, start_left + Inches(0.2), top + Inches(0.3), col_width - Inches(0.4), height - Inches(0.6), [
        "• YOLO26n 自訂模型：使用自行訓練的輕量化模型，進行即時影像中人員的安全帽及反光背心穿戴偵測。",
        "• 毫秒級邊緣運算：支援網頁 Web 相機或圖片串流，在邊緣端快速完成影像辨識，確保高效率偵測。",
        "• 精準比例座標標記：將檢測框轉換為相對比例，保障網頁 Canvas 縮放時的精準渲染。"
    ], title="📷 邊緣端影像監控", title_size=18, text_size=13)
    
    # 卡片 2 (Discord)
    create_card(slide, start_left + col_width + gap, top, col_width, height, RGBColor(255, 255, 255), RGBColor(37, 99, 235))
    add_bullet_points(slide, start_left + col_width + gap + Inches(0.2), top + Inches(0.3), col_width - Inches(0.4), height - Inches(0.6), [
        "• 雙階段警報機制：違規當下立即發送警告文字與畫面截圖；影片預錄完成後，自動非同步推送回放影片至 Discord。",
        "• 違規影片事證：直接在 Discord 接收 10 秒違規影片檔，省去開啟網頁系統的時間，提供主管第一手影像事證。",
        "• 去抖動與冷卻機制：內建 10 秒警報冷卻時間，避免在連續畫面中重複發送警報而洗版。"
    ], title="🚨 Discord 雙階即時警報", title_size=18, text_size=13)
    
    # 卡片 3 (Gemini)
    create_card(slide, start_left + (col_width + gap)*2, top, col_width, height, RGBColor(255, 255, 255), RGBColor(139, 92, 246))
    add_bullet_points(slide, start_left + (col_width + gap)*2 + Inches(0.2), top + Inches(0.3), col_width - Inches(0.4), height - Inches(0.6), [
        "• 智慧稽核分析：整合 Google Gemini 3.5 Flash API，在「Gemini 智慧稽核」分頁提供一鍵生成分析功能。",
        "• 大數據統計：自動撈取資料庫中最新的 100 筆違規紀錄，由 AI 進行違規統計與趨勢分析。",
        "• 管理改善建議：由 AI 評估工安風險，並產出繁體中文的精準改善與現場防範建議報告。"
    ], title="🧠 Gemini AI 智慧報告", title_size=18, text_size=13)
    
    # ----------------------------------------------------
    # Slide 4: 系統架構設計與模組關係
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, RGBColor(248, 250, 252))
    create_title(slide, "03", "系統架構設計與模組關係")
    
    row_height = Inches(1.3)
    row_gap = Inches(0.35)
    start_top = Inches(1.8)
    row_width = Inches(11.7)
    left = Inches(0.8)
    
    # Layer 1
    create_card(slide, left, start_top, row_width, row_height, RGBColor(255, 255, 255), RGBColor(16, 185, 129))
    add_bullet_points(slide, left + Inches(0.2), start_top + Inches(0.15), row_width - Inches(0.4), row_height - Inches(0.2), [
        "• 使用背景 Python YOLO Worker：載入自訂訓練的 `yolo26n.pt` 權重，進行高效的即時畫面物體偵測。",
        "• 雙向進程通訊：與後端 Node.js 通過 IPC (stdin/stdout JSON) 管道，進行毫秒級的物體檢測數據對接。"
    ], title="1. 影像擷取與邊緣 AI 推論層", title_size=16, text_size=13)
    
    # Layer 2
    create_card(slide, left, start_top + row_height + row_gap, row_width, row_height, RGBColor(255, 255, 255), RGBColor(37, 99, 235))
    add_bullet_points(slide, left + Inches(0.2), start_top + row_height + row_gap + Inches(0.15), row_width - Inches(0.4), row_height - Inches(0.2), [
        "• Node.js Express 後端伺服器：作為系統的控制核心，負責協調前端、YOLO Worker、資料庫與外部 API。",
        "• 數據持久化與雙向通訊：使用 Sequelize ORM 連接 PostgreSQL 資料庫儲存違規事件；透過 Socket.io 實時推送日誌至前端。"
    ], title="2. 核心控制與數據儲存層 (Node.js & PostgreSQL)", title_size=16, text_size=13)
    
    # Layer 3
    create_card(slide, left, start_top + (row_height + row_gap)*2, row_width, row_height, RGBColor(255, 255, 255), RGBColor(139, 92, 246))
    add_bullet_points(slide, left + Inches(0.2), start_top + (row_height + row_gap)*2 + Inches(0.15), row_width - Inches(0.4), row_height - Inches(0.2), [
        "• Google Gemini 3.5 Flash API：負責調用大語言模型，對 PostgreSQL 中的違規數據進行深度稽核與分析報告生成。",
        "• Discord Webhook 警報：負責向外部 Discord 頻道發送包含即時違規影像截圖的訊息通知，並於影片上傳後自動推送影片回放。"
    ], title="3. 智慧分析與外部警報外接層 (Gemini & Discord)", title_size=16, text_size=13)
    
    # ----------------------------------------------------
    # Slide 5: 智慧安全偵測與通報流程
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, RGBColor(248, 250, 252))
    create_title(slide, "04", "智慧安全偵測與通報流程")
    
    step_width = Inches(2.6)
    arrow_width = Inches(0.4)
    step_gap = Inches(0.1)
    start_left = Inches(0.8)
    top = Inches(2.2)
    height = Inches(4.0)
    
    steps_data = [
        {"num": "1", "title": "影像擷取與傳輸", "desc": ["• 網頁前端調用 Web 相機實時擷取畫面影格", "• 透過影像串流或 Base64 格式即時傳送至後端進行處理"]},
        {"num": "2", "title": "YOLO 智慧偵測", "desc": ["• 後端載入 `yolo26n.pt` 進行影像物體偵測", "• 毫秒級判定人員是否戴安全帽與穿反光背心"]},
        {"num": "3", "title": "違規事件存庫", "desc": ["• 當判定有違規行為時，自動截圖並以時間戳命名", "• 將違規時間、類別與截圖路徑寫入 PostgreSQL"]},
        {"num": "4", "title": "雙階 Discord 通報", "desc": ["• 階段一：實時發送警報文字與現場截圖至 Discord", "• 階段二：5 秒後前端合成預錄影片上傳，後端非同步推送影片"]}
    ]
    
    for i, step in enumerate(steps_data):
        x = start_left + i * (step_width + arrow_width + step_gap)
        # 畫步驟卡片
        create_card(slide, x, top, step_width, height, RGBColor(255, 255, 255), RGBColor(37, 99, 235))
        
        # 步驟數字標籤 (小圓形)
        num_box = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), top + Inches(0.2), Inches(0.6), Inches(0.6))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = RGBColor(37, 99, 235)
        num_box.line.fill.background()
        num_tf = num_box.text_frame
        num_tf.margin_left = num_tf.margin_top = num_tf.margin_right = num_tf.margin_bottom = 0
        p = num_tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = step["num"]
        p.font.name = 'Arial'
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # 步驟標題與內容
        add_bullet_points(slide, x + Inches(0.1), top + Inches(0.9), step_width - Inches(0.2), height - Inches(1.0), 
                          step["desc"], title=step["title"], title_size=15, text_size=12)
        
        # 畫箭頭 (除了最後一個)
        if i < 3:
            arrow_x = x + step_width + Inches(0.05)
            arrow_y = top + Inches(1.8)
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, arrow_y, arrow_width, Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(148, 163, 184) # Slate-400
            arrow.line.fill.background()
            
    # ----------------------------------------------------
    # Slide 6: Gemini AI 智慧工安報告
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, RGBColor(248, 250, 252))
    create_title(slide, "05", "Gemini AI 智慧工安報告與稽核")
    
    # 左欄
    create_card(slide, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0), RGBColor(255, 255, 255), RGBColor(139, 92, 246))
    add_bullet_points(slide, Inches(1.1), Inches(1.9), Inches(5.0), Inches(4.4), [
        "• 一鍵自動化分析：管理人員只需在「Gemini 智慧稽核」分頁點擊按鈕，系統即自動調用 API，無須手動整理數據。",
        "• 資料庫無縫對接：後端會自動撈取 PostgreSQL 中最新的 100 筆違規紀錄（包含時間、違規類型等數據）。",
        "• 繁體中文智慧生成：AI 專門針對台灣工安規範，以繁體中文生成排版工整、論點清晰的分析與建議報告。"
    ], title="AI 智慧稽核運作模式", title_size=20, text_size=15)
    
    # 右欄
    create_card(slide, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.0), RGBColor(255, 255, 255), RGBColor(139, 92, 246))
    add_bullet_points(slide, Inches(7.2), Inches(1.9), Inches(5.0), Inches(4.4), [
        "• 違規統計分析：由 Gemini 自動歸納出最常發生違規的時段（例如交班或疲勞時段）以及最常被忽略的防護具（如安全帽或背心）。",
        "• 工安風險評估：依據違規頻率與型態，量化廠區的安全風險指數，協助管理者快速掌握當前防線的脆弱環節。",
        "• 管理改善建議：提供精準且可落地的防範對策，例如調整現場宣導、加強特定區域巡檢，或優化人員動線。"
    ], title="報告三大核心價值", title_size=20, text_size=15)

    # ----------------------------------------------------
    # Slide 7: 違規前後預錄與歷史回溯
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, RGBColor(248, 250, 252))
    create_title(slide, "06", "違規前後預錄與歷史回溯")
    
    # 左欄 (前端錄製機制)
    create_card(slide, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0), RGBColor(255, 255, 255), RGBColor(37, 99, 235))
    add_bullet_points(slide, Inches(1.1), Inches(1.9), Inches(5.0), Inches(4.4), [
        "• HTML5 Canvas 錄製：前端利用 `MediaRecorder` API 監聽 `captureStream(10)`，實現零伺服器效能開銷的錄影。",
        "• 5+5 雙向預錄機制：前端隨時維持 5 秒的環形預錄緩衝區；當違規觸發時，繼續錄製 5 秒，隨後合併上傳。",
        "• 保留 YOLO 偵測框線：直接錄製前端 Canvas 畫面，使得回放影片中可完整保留 YOLO 的紅色違規框標籤。"
    ], title="前端 Canvas 雙向預錄", title_size=20, text_size=15)
    
    # 右欄 (後端整合與 UI 播放)
    create_card(slide, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.0), RGBColor(255, 255, 255), RGBColor(37, 99, 235))
    add_bullet_points(slide, Inches(7.2), Inches(1.9), Inches(5.0), Inches(4.4), [
        "• 後端二進位儲存：API 接收二進位影片後，將影片寫入 `public/videos/`，並自動更新資料庫的 `videoPath` 欄位。",
        "• 動態回放按鈕：影片上傳後，後端透過 Socket.io 廣播 `video_ready`，前端日誌列表將動態出現「🎥 回放」按鈕。",
        "• 彈出式質感播放器：點擊播放按鈕，在畫面上跳出精美的半透明磨砂玻璃質感 Modal，並流暢播放 WebM 影片。"
    ], title="後端影片關聯與 Modal 播放", title_size=20, text_size=15)

    # ----------------------------------------------------
    # Slide 8: 系統人性化互動設計 (電子圍籬)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, RGBColor(248, 250, 252))
    create_title(slide, "07", "系統人性化互動設計")
    
    # 左欄 (電子圍籬與幾何判定)
    create_card(slide, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0), RGBColor(255, 255, 255), RGBColor(239, 68, 68))
    add_bullet_points(slide, Inches(1.1), Inches(1.9), Inches(5.0), Inches(4.4), [
        "• 多邊形圍籬繪製：使用者點擊 Canvas 即可繪製任意形狀的多邊形管制區，並提供半透明紅色的區域視覺提示。",
        "• 座標相對比例化：所有頂點座標以 0.0 ~ 1.0 的相對比例儲存，確保在瀏覽器縮放與不同解析度下百分之百精準。",
        "• Point-in-Polygon 判定：後端在 Node.js 中實作高效的射線法，計算人員底邊中心點，非管制區內的违規將自動過濾。"
    ], title="📐 電子圍籬與空間判定", title_size=20, text_size=15)
    
    # 右欄 (無串流編輯與防禦性優化)
    create_card(slide, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.0), RGBColor(255, 255, 255), RGBColor(239, 68, 68))
    add_bullet_points(slide, Inches(7.2), Inches(1.9), Inches(5.0), Inches(4.4), [
        "• 無串流靜態編輯與預覽：支援在未啟動相機/IP Cam 時，依然能繪製、清除並載入預設圍籬，提供靜態 Canvas 渲染。",
        "• 滑鼠穿透控制優化：進入編輯狀態時自動修改為 `pointer-events: auto`，儲存後恢復 `none` 避免干擾按鈕操作。",
        "• 防禦性座標設計：避免 Canvas 寬高在初始化為 0 時計算出 NaN 座標，增加防錯除錯，保障代碼的健壯度。"
    ], title="⚙️ 人性化交互與防錯設計", title_size=20, text_size=15)

    # ----------------------------------------------------
    # Slide 9: 結論與未來展望
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, RGBColor(248, 250, 252))
    create_title(slide, "08", "總結與未來展望")
    
    # 左欄
    create_card(slide, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0), RGBColor(255, 255, 255), RGBColor(16, 185, 129))
    add_bullet_points(slide, Inches(1.1), Inches(1.9), Inches(5.0), Inches(4.4), [
        "• 被動防護轉為主動防禦：跳脫傳統事後調閱監視器的模式，在違規發生的第一秒立即通報，將意外機率降至最低。",
        "• 工安管理數位化：結合資料庫與 AI 分析，將冰冷的監控影像轉化為具備管理價值的工安稽核報告與影片事證庫。",
        "• 顯著降低巡檢人力：全天候自動化監控，減輕安全稽核人員現場巡查負擔，提升管理效率。"
    ], title="🏆 系統實施成效", title_size=20, text_size=15)
    
    # 右欄
    create_card(slide, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.0), RGBColor(255, 255, 255), RGBColor(37, 99, 235))
    add_bullet_points(slide, Inches(7.2), Inches(1.9), Inches(5.0), Inches(4.4), [
        "• 擴充更多防護具識別：未來計劃引入安全鞋、防墜安全帶、護目鏡或口罩等多種類別的 YOLO 偵測模型。",
        "• 多相機與多廠區支援：支援接入多路 IP Camera，實現大規模廠區或跨廠區的同步邊緣監控與集中管理。",
        "• 雙平台通報機制：除目前已實現的 Discord 通報外，未來可啟用並完善 LINE Notify 通報，滿足多元通訊需求。"
    ], title="🚀 未來擴充方向", title_size=20, text_size=15)
    
    # ----------------------------------------------------
    # Slide 10: 封底
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, RGBColor(15, 23, 42)) # Slate 900
    
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.33), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = 0
    
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    p1.text = "簡報結束，謝謝聆聽"
    p1.font.name = 'Microsoft JhengHei'
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)
    p1.space_after = Pt(20)
    
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = "AIoT 智慧安全監控與主動防禦平台"
    p2.font.name = 'Microsoft JhengHei'
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(148, 163, 184)
    
    # 保存
    output_filename = "AIoT_Safety_Platform_Report.pptx"
    prs.save(output_filename)
    print(f"Successfully generated {output_filename}")

if __name__ == '__main__':
    main()
